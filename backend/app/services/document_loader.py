# backend/app/services/document_loader.py
import os
from typing import List

class DocumentLoader:
    """Extract text from PDF, PPT, Word, and TXT files"""
    
    @staticmethod
    def load_pdf(file_path: str) -> str:
        try:
            from PyPDF2 import PdfReader
            text = ""
            reader = PdfReader(file_path)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            return text
        except ImportError:
            print("PyPDF2 not installed. Run: pip install pypdf2")
            return ""
    
    @staticmethod
    def load_docx(file_path: str) -> str:
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(file_path)
            return "\n".join([p.text for p in doc.paragraphs if p.text])
        except ImportError:
            print("python-docx not installed. Run: pip install python-docx")
            return ""
    
    @staticmethod
    def load_pptx(file_path: str) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
            return text
        except ImportError:
            print("python-pptx not installed. Run: pip install python-pptx")
            return ""
    
    @staticmethod
    def load_txt(file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    @staticmethod
    def load_file(file_path: str) -> str:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.pdf':
            return DocumentLoader.load_pdf(file_path)
        elif ext in ['.docx', '.doc']:
            return DocumentLoader.load_docx(file_path)
        elif ext in ['.pptx', '.ppt']:
            return DocumentLoader.load_pptx(file_path)
        elif ext == '.txt':
            return DocumentLoader.load_txt(file_path)
        else:
            print(f"Unsupported file type: {ext}")
            return ""
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = 300) -> List[str]:
        words = text.split()
        chunks = []
        for i in range(0, len(words), chunk_size):
            chunk = " ".join(words[i:i+chunk_size])
            if len(chunk) > 20:
                chunks.append(chunk)
        return chunks
    
    @staticmethod
    def load_folder(folder_path: str) -> List[str]:
        all_chunks = []
        if not os.path.exists(folder_path):
            print(f"Folder not found: {folder_path}")
            return all_chunks
        
        for filename in os.listdir(folder_path):
            file_path = os.path.join(folder_path, filename)
            if os.path.isfile(file_path):
                try:
                    text = DocumentLoader.load_file(file_path)
                    if text:
                        chunks = DocumentLoader.chunk_text(text)
                        all_chunks.extend(chunks)
                        print(f"  Loaded: {filename} -> {len(chunks)} chunks")
                except Exception as e:
                    print(f"  Failed: {filename} - {e}")
        return all_chunks

document_loader = DocumentLoader()
